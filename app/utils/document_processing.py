from PyPDF2 import PdfReader
from pptx import Presentation
from langchain.schema import Document

def create_document(file):
    filename = file.filename.lower()
    if filename.endswith('.pdf'):
        pdf_reader = PdfReader(file)
        text = "".join(page.extract_text() for page in pdf_reader.pages)
    elif filename.endswith('.pptx'):
        prs = Presentation(file)
        text_list = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_list.append(shape.text)
        text = "\n".join(text_list)
    else:
        raise ValueError("Unsupported file format. Only PDF and PPTX are supported.")

    # Filter content: remove blank lines
    text = "\n".join(line for line in text.split("\n") if line.strip())

    # Create a Document instance with the full text
    document = Document(page_content=text, metadata={})
    return [document]
